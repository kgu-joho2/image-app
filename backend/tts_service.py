import os
import io
import base64
import tempfile
import traceback
import wave
from typing import Dict, Any, Optional, Union
from google import genai
from google.genai import types
from google.cloud import documentai
from google.cloud import vision
import PyPDF2
import docx
from pptx import Presentation
from PIL import Image
import speech_recognition as sr

def wave_file(pcm_data, channels=1, rate=24000, sample_width=2):
    """
    PCMデータをWAVファイル形式に変換
    """
    wav_buffer = io.BytesIO()
    with wave.open(wav_buffer, "wb") as wf:
        wf.setnchannels(channels)
        wf.setsampwidth(sample_width)
        wf.setframerate(rate)
        wf.writeframes(pcm_data)
    wav_buffer.seek(0)
    return wav_buffer.getvalue()

class TTSService:
    def __init__(self):
        # Configure Gemini API
        api_key = os.getenv("GOOGLE_API_KEY")
        if not api_key:
            raise ValueError("GOOGLE_API_KEY environment variable not set.")
        
        # Initialize Gemini client
        self.client = genai.Client(api_key=api_key)
        
        # Initialize Document AI client if available
        self.document_ai_client = None
        try:
            if os.getenv("GOOGLE_APPLICATION_CREDENTIALS"):
                self.document_ai_client = documentai.DocumentProcessorServiceClient()
                self.project_id = os.getenv("GOOGLE_CLOUD_PROJECT_ID")
                self.processor_id = os.getenv("DOCUMENT_AI_PROCESSOR_ID")
                self.location = os.getenv("DOCUMENT_AI_LOCATION", "us")
        except Exception as e:
            print(f"Document AI initialization failed: {e}")
        
        # Initialize Vision API client if available
        self.vision_client = None
        try:
            if os.getenv("GOOGLE_APPLICATION_CREDENTIALS"):
                self.vision_client = vision.ImageAnnotatorClient()
        except Exception as e:
            print(f"Vision API initialization failed: {e}")

    def extract_text_from_file(self, file_content: bytes, file_type: str, filename: str) -> Dict[str, Any]:
        """
        ファイルからテキストを抽出します
        """
        try:
            if file_type == 'application/pdf':
                return self._extract_from_pdf(file_content)
            elif file_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
                return self._extract_from_docx(file_content)
            elif file_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
                return self._extract_from_pptx(file_content)
            elif file_type == 'text/plain':
                return self._extract_from_txt(file_content)
            elif file_type in ['image/jpeg', 'image/png']:
                return self._extract_from_image(file_content)
            else:
                raise ValueError(f"Unsupported file type: {file_type}")
                
        except Exception as e:
            print(f"Text extraction error: {e}")
            traceback.print_exc()
            return {
                "success": False,
                "error": f"テキスト抽出に失敗しました: {str(e)}"
            }

    def _extract_from_pdf(self, file_content: bytes) -> Dict[str, Any]:
        """PDFファイルからテキストを抽出"""
        try:
            # Document AI使用を優先
            if self.document_ai_client and self.project_id and self.processor_id:
                return self._extract_with_document_ai(file_content, "application/pdf")
            
            # フォールバック: PyPDF2を使用
            text = ""
            pdf_file = io.BytesIO(file_content)
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
            
            return {
                "success": True,
                "text": text.strip(),
                "method": "PyPDF2"
            }
            
        except Exception as e:
            return {
                "success": False,
                "error": f"PDF処理エラー: {str(e)}"
            }

    def _extract_from_docx(self, file_content: bytes) -> Dict[str, Any]:
        """DOCXファイルからテキストを抽出"""
        try:
            doc_file = io.BytesIO(file_content)
            doc = docx.Document(doc_file)
            
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            
            # テーブルからもテキストを抽出
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text += cell.text + " "
                    text += "\n"
            
            return {
                "success": True,
                "text": text.strip(),
                "method": "python-docx"
            }
            
        except Exception as e:
            return {
                "success": False,
                "error": f"DOCX処理エラー: {str(e)}"
            }

    def _extract_from_pptx(self, file_content: bytes) -> Dict[str, Any]:
        """PPTXファイルからテキストを抽出"""
        try:
            ppt_file = io.BytesIO(file_content)
            presentation = Presentation(ppt_file)
            
            text = ""
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            
            return {
                "success": True,
                "text": text.strip(),
                "method": "python-pptx"
            }
            
        except Exception as e:
            return {
                "success": False,
                "error": f"PPTX処理エラー: {str(e)}"
            }

    def _extract_from_txt(self, file_content: bytes) -> Dict[str, Any]:
        """TXTファイルからテキストを抽出"""
        try:
            # UTF-8でデコードを試行
            try:
                text = file_content.decode('utf-8')
            except UnicodeDecodeError:
                # フォールバック: CP932やShift_JISを試行
                try:
                    text = file_content.decode('cp932')
                except UnicodeDecodeError:
                    text = file_content.decode('shift_jis')
            
            return {
                "success": True,
                "text": text.strip(),
                "method": "text-decode"
            }
            
        except Exception as e:
            return {
                "success": False,
                "error": f"TXT処理エラー: {str(e)}"
            }

    def _extract_from_image(self, file_content: bytes) -> Dict[str, Any]:
        """画像ファイルからOCRでテキストを抽出"""
        try:
            # Vision API使用を優先
            if self.vision_client:
                return self._extract_with_vision_api(file_content)
            
            # フォールバック（基本OCRライブラリがあれば）
            return {
                "success": False,
                "error": "OCR機能が利用できません。Google Cloud Vision APIの設定が必要です。"
            }
            
        except Exception as e:
            return {
                "success": False,
                "error": f"画像処理エラー: {str(e)}"
            }

    def _extract_with_document_ai(self, file_content: bytes, mime_type: str) -> Dict[str, Any]:
        """Google Cloud Document AIを使用してテキストを抽出"""
        try:
            # Document AI processor name
            name = f"projects/{self.project_id}/locations/{self.location}/processors/{self.processor_id}"
            
            # Raw document
            raw_document = documentai.RawDocument(content=file_content, mime_type=mime_type)
            
            # Process request
            request = documentai.ProcessRequest(name=name, raw_document=raw_document)
            result = self.document_ai_client.process_document(request=request)
            
            # Extract text
            text = result.document.text
            
            return {
                "success": True,
                "text": text.strip(),
                "method": "Document AI"
            }
            
        except Exception as e:
            return {
                "success": False,
                "error": f"Document AI処理エラー: {str(e)}"
            }

    def _extract_with_vision_api(self, file_content: bytes) -> Dict[str, Any]:
        """Google Cloud Vision APIを使用してOCR"""
        try:
            image = vision.Image(content=file_content)
            response = self.vision_client.text_detection(image=image)
            
            if response.error.message:
                raise Exception(f"Vision API error: {response.error.message}")
            
            texts = response.text_annotations
            if texts:
                text = texts[0].description
                return {
                    "success": True,
                    "text": text.strip(),
                    "method": "Vision API OCR"
                }
            else:
                return {
                    "success": True,
                    "text": "",
                    "method": "Vision API OCR",
                    "note": "テキストが検出されませんでした"
                }
                
        except Exception as e:
            return {
                "success": False,
                "error": f"Vision API処理エラー: {str(e)}"
            }

    def summarize_text(self, text: str, speaker_mode: str = "single") -> Dict[str, Any]:
        """
        テキストを要約します（話者モードに応じて）
        """
        try:
            # テキスト要約用には従来のGenerativeModelを使用
            import google.generativeai as genai_classic
            genai_classic.configure(api_key=os.getenv("GOOGLE_API_KEY"))
            model = genai_classic.GenerativeModel("gemini-2.0-flash")
            
            if speaker_mode == "single":
                # 単一話者用要約
                prompt = f"""
                以下の文書の内容を、音声読み上げに適した形で要約してください。
                
                要約の条件:
                1. 重要なポイントを漏らさずに簡潔にまとめる
                2. 音声で聞いた時に理解しやすい構造にする
                3. 専門用語は必要に応じて説明を加える
                4. 自然な日本語で、読み上げに適した文体にする
                5. 一人の話者が読み上げる形式で要約する
                6. 要約内容のみを出力し、説明文や前置きは不要
                7. 「要約します」「以下のような内容です」等の文言は含めない
                
                文書内容:
                {text}
                """
            else:
                # 複数話者用要約（会話形式）
                prompt = f"""
                以下の文書の内容を、2人の話者（話者A、話者B）が会話する形式で要約してください。
                
                要約の条件:
                1. 文書の重要なポイントを会話形式で分かりやすく説明
                2. 話者A と 話者B が交互に話す自然な会話形式
                3. 専門用語や重要な概念は会話の中で説明
                4. 各発言は1〜2文程度で簡潔に
                5. 文書の内容を正確に反映した会話内容
                6. 音声読み上げに適した自然な日本語
                
                出力形式:
                話者A: [発言内容]
                話者B: [発言内容]
                話者A: [発言内容]
                ...
                
                文書内容:
                {text}
                """
            
            response = model.generate_content(prompt)
            
            if (response.candidates and 
                response.candidates[0].content and 
                response.candidates[0].content.parts and 
                response.candidates[0].content.parts[0].text):
                
                summary = response.candidates[0].content.parts[0].text.strip()
                
                return {
                    "success": True,
                    "summary": summary,
                    "speaker_mode": speaker_mode
                }
            else:
                return {
                    "success": False,
                    "error": "要約の生成に失敗しました"
                }
                
        except Exception as e:
            print(f"Summarization error: {e}")
            traceback.print_exc()
            return {
                "success": False,
                "error": f"要約生成エラー: {str(e)}"
            }

    def generate_speech(self, text: str, voice_settings: Dict[str, Any], 
                       speaker_mode: str = "single", style: str = "", 
                       rate: float = 1.0) -> Union[bytes, Dict[str, Any]]:
        """
        正しいGemini 2.5 TTS APIを使用してテキストから音声を生成します
        """
        try:
            print(f"TTS生成開始: speaker_mode={speaker_mode}, voice_settings={voice_settings}")
            print(f"テキスト: '{text}', スタイル: '{style}', レート: {rate}")
            
            # 音声生成の設定
            if speaker_mode == "single":
                # 単一話者モード
                voice_name = voice_settings.get('voice', 'Kore')
                
                # プロンプトにスタイル指示を含める
                prompt = text
                if style:
                    prompt = f"音声スタイル: {style}\n\n{text}"
                
                print(f"単一話者モード: voice={voice_name}, prompt='{prompt}'")
                
                try:
                    response = self.client.models.generate_content(
                        model="gemini-2.5-flash-preview-tts",
                        contents=prompt,
                        config=types.GenerateContentConfig(
                            response_modalities=["AUDIO"],
                            speech_config=types.SpeechConfig(
                                voice_config=types.VoiceConfig(
                                    prebuilt_voice_config=types.PrebuiltVoiceConfig(
                                        voice_name=voice_name
                                    )
                                )
                            )
                        )
                    )
                    print(f"単一話者 API応答受信: {type(response)}")
                except Exception as api_error:
                    print(f"単一話者 API呼び出しエラー: {api_error}")
                    # 500エラーの場合は再試行
                    if "500" in str(api_error) or "INTERNAL" in str(api_error):
                        print("500エラー検出、3秒後に再試行...")
                        import time
                        time.sleep(3)
                        try:
                            response = self.client.models.generate_content(
                                model="gemini-2.5-flash-preview-tts",
                                contents=prompt,
                                config=types.GenerateContentConfig(
                                    response_modalities=["AUDIO"],
                                    speech_config=types.SpeechConfig(
                                        voice_config=types.VoiceConfig(
                                            prebuilt_voice_config=types.PrebuiltVoiceConfig(
                                                voice_name=voice_name
                                            )
                                        )
                                    )
                                )
                            )
                            print("再試行成功")
                        except Exception as retry_error:
                            print(f"再試行も失敗: {retry_error}")
                            raise retry_error
                    else:
                        raise api_error
                        
            else:
                # 複数話者モード
                voice_a = voice_settings.get('voiceA', 'Kore')
                voice_b = voice_settings.get('voiceB', 'Puck')
                
                # 複数話者用のプロンプト - フォーマットを改善
                speaker_text = text
                if not any(keyword in text for keyword in [':', '：', '話者', 'Speaker']):
                    # 話者指定がない場合は自動的に分割
                    lines = text.split('\n')
                    formatted_lines = []
                    for i, line in enumerate(lines):
                        if line.strip():
                            speaker = '話者A' if i % 2 == 0 else '話者B'
                            formatted_lines.append(f"{speaker}: {line.strip()}")
                    speaker_text = '\n'.join(formatted_lines)
                
                prompt = f"以下の会話を2人の話者で読み上げてください:\n{speaker_text}"
                if style:
                    prompt = f"音声スタイル: {style}\n\n{prompt}"
                
                print(f"複数話者モード: voiceA={voice_a}, voiceB={voice_b}")
                print(f"複数話者プロンプト: '{prompt}'")
                
                try:
                    response = self.client.models.generate_content(
                        model="gemini-2.5-flash-preview-tts",
                        contents=prompt,
                        config=types.GenerateContentConfig(
                            response_modalities=["AUDIO"],
                            speech_config=types.SpeechConfig(
                                multi_speaker_voice_config=types.MultiSpeakerVoiceConfig(
                                    speaker_voice_configs=[
                                        types.SpeakerVoiceConfig(
                                            speaker='話者A',
                                            voice_config=types.VoiceConfig(
                                                prebuilt_voice_config=types.PrebuiltVoiceConfig(
                                                    voice_name=voice_a
                                                )
                                            )
                                        ),
                                        types.SpeakerVoiceConfig(
                                            speaker='話者B',
                                            voice_config=types.VoiceConfig(
                                                prebuilt_voice_config=types.PrebuiltVoiceConfig(
                                                    voice_name=voice_b
                                                )
                                            )
                                        )
                                    ]
                                )
                            )
                        )
                    )
                    print(f"複数話者 API応答受信: {type(response)}")
                except Exception as api_error:
                    print(f"複数話者 API呼び出しエラー: {api_error}")
                    raise api_error
            
            # 音声データを取得
            print("レスポンス解析開始...")
            print(f"レスポンス全体: {response}")
            
            if response and hasattr(response, 'candidates') and response.candidates:
                print(f"candidates数: {len(response.candidates)}")
                candidate = response.candidates[0]
                print(f"candidate: {candidate}")
                
                if hasattr(candidate, 'content') and candidate.content:
                    print(f"content: {candidate.content}")
                    print(f"content parts数: {len(candidate.content.parts) if candidate.content.parts else 0}")
                    
                    if candidate.content.parts and len(candidate.content.parts) > 0:
                        part = candidate.content.parts[0]
                        print(f"part: {part}")
                        print(f"part type: {type(part)}")
                        print(f"part attributes: {dir(part)}")
                        
                        if hasattr(part, 'inline_data') and part.inline_data:
                            print(f"inline_data found: {type(part.inline_data)}")
                            print(f"inline_data attributes: {dir(part.inline_data)}")
                            
                            # PCMデータを取得
                            pcm_data = part.inline_data.data
                            print(f"PCMデータサイズ: {len(pcm_data)} bytes")
                            
                            # PCMデータをWAVファイル形式に変換
                            wav_data = wave_file(pcm_data)
                            print(f"WAVデータサイズ: {len(wav_data)} bytes")
                            
                            # base64エンコードして返す
                            wav_data_b64 = base64.b64encode(wav_data).decode('utf-8')
                            
                            return {
                                "success": True,
                                "audio_data": wav_data_b64,
                                "format": "wav"
                            }
                        elif hasattr(part, 'text') and part.text:
                            print(f"テキストレスポンス発見: {part.text}")
                            print("音声データではなくテキストレスポンスが返されました")
                        else:
                            print("inline_data が見つかりません")
                            print(f"part 内容: {part}")
                    else:
                        print("parts が空またはありません")
                else:
                    print("content が見つかりません")
            else:
                print("candidates が見つかりません")
            
            print("音声データの取得に失敗: 期待されるデータ構造が見つかりませんでした")
            return {
                "success": False,
                "error": "音声データの生成に失敗しました - レスポンス構造が不正です"
            }
                
        except Exception as e:
            print(f"Speech generation error: {e}")
            traceback.print_exc()
            return {
                "success": False,
                "error": f"音声生成エラー: {str(e)}"
            }

    def preview_voice(self, voice: str, text: str = "こんにちは。これは音声のプレビューです。", 
                     style: str = "", rate: float = 1.0) -> Union[bytes, Dict[str, Any]]:
        """
        音声のプレビューを生成します
        """
        try:
            # プレビュー用のテキストで音声生成
            voice_settings = {'voice': voice}
            result = self.generate_speech(text, voice_settings, "single", style, rate)
            
            if result.get("success"):
                return {
                    "success": True,
                    "audio_data": result["audio_data"],
                    "format": "wav"
                }
            else:
                return result
            
        except Exception as e:
            print(f"Voice preview error: {e}")
            return {
                "success": False,
                "error": f"音声プレビューエラー: {str(e)}"
            } 